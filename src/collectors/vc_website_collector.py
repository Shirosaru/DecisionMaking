from __future__ import annotations

"""
VC Website Slide & Document Collector
======================================
Comprehensive scraper for investor presentations, pitch decks, and pipeline
slides from top-tier bio VC firms, biotech IR pages, and healthcare conference
archives.  Goal: learn GO / NO-GO decision-making patterns from the best.

Strategy per VC site:
  1. Fetch the firm's news/insights/press or portfolio page.
  2. Find all links to PDF / PPTX files or slide-hosting services
     (SlideShare, SpeakerDeck, IR pages, JPM, Cowen, Jefferies CDNs …).
  3. For PDF/PPTX links, download directly → data/slides/<sub>/<vc_slug>/.
  4. For HTML pages that host slide content, scrape the text.
  5. Run NLP to extract clinical-stage / decision / indication / mechanism.

Sub-directories:
  data/slides/vc/          — VC firm news & portfolio pages
  data/slides/conference/  — Healthcare conference archives
  data/slides/startup/     — Startup / biotech company IR presentations
  data/slides/edgar/       — SEC EDGAR EX-99 exhibits (slide_downloader.py)

PDF extraction requires pdfplumber; PPTX extraction requires python-pptx.
Both are listed in requirements.txt.
"""

import hashlib
import io
import logging
import re
from pathlib import Path
from typing import Any
from urllib.parse import urljoin, urlparse

from bs4 import BeautifulSoup

from .base_collector import BaseCollector, RawRecord

logger = logging.getLogger(__name__)

# ── Source catalogue ──────────────────────────────────────────────────────────
#
# Fields:
#   vc      — display name / grouping label
#   url     — seed URL to fetch
#   depth   — 0 = seed page only; 1 = follow PDF/slide links one level deep
#   rss     — True if URL is an RSS/Atom feed
#   subdir  — override output sub-directory (default: "vc")
#
# ─────────────────────────────────────────────────────────────────────────────

_VC_PAGES: list[dict[str, Any]] = [

    # ═══════════════════════════════════════════════════════════════════════
    # TIER-1 BIOTECH VC FIRMS
    # ═══════════════════════════════════════════════════════════════════════

    # Flagship Pioneering
    {"vc": "Flagship Pioneering",       "url": "https://www.flagshippioneering.com/news",                     "depth": 1},
    {"vc": "Flagship Pioneering",       "url": "https://www.flagshippioneering.com/companies",                "depth": 1},
    # Atlas Venture
    {"vc": "Atlas Venture",             "url": "https://atlasventure.com/news/",                             "depth": 1},
    {"vc": "Atlas Venture",             "url": "https://lifescivc.com",                                      "depth": 1},
    # Third Rock Ventures
    {"vc": "Third Rock Ventures",       "url": "https://thirdrockventures.com/news",                          "depth": 1},
    {"vc": "Third Rock Ventures",       "url": "https://thirdrockventures.com/portfolio",                     "depth": 1},
    # ARCH Venture Partners
    {"vc": "ARCH Venture Partners",     "url": "https://www.archventure.com/news",                           "depth": 1},
    {"vc": "ARCH Venture Partners",     "url": "https://medium.com/arch-venture-partners",                   "depth": 1},
    # Foresite Capital
    {"vc": "Foresite Capital",          "url": "https://foresitecapital.com/insights",                       "depth": 1},
    {"vc": "Foresite Capital",          "url": "https://foresitecapital.com/portfolio",                      "depth": 1},
    # Versant Ventures
    {"vc": "Versant Ventures",          "url": "https://www.versantventures.com/news",                       "depth": 1},
    {"vc": "Versant Ventures",          "url": "https://www.versantventures.com/portfolio",                  "depth": 0},
    # RA Capital Management
    {"vc": "RA Capital",                "url": "https://www.racap.com/news",                                 "depth": 1},
    {"vc": "RA Capital",                "url": "https://www.racap.com/perspectives",                         "depth": 1},
    # OrbiMed Advisors
    {"vc": "OrbiMed",                   "url": "https://www.orbimed.com/news",                               "depth": 1},
    # 5AM Ventures
    {"vc": "5AM Ventures",              "url": "https://www.5amventures.com/news",                           "depth": 1},
    # Novo Holdings Equity
    {"vc": "Novo Holdings",             "url": "https://www.novoholdings.com/news",                          "depth": 1},
    {"vc": "Novo Holdings",             "url": "https://www.novoholdings.com/investments/life-sciences",     "depth": 1},
    # Sofinnova Partners (EU)
    {"vc": "Sofinnova Partners",        "url": "https://www.sofinnova.com/news",                             "depth": 1},
    # Omega Funds
    {"vc": "Omega Funds",               "url": "https://www.omegafunds.com/news",                            "depth": 1},
    # Vida Ventures
    {"vc": "Vida Ventures",             "url": "https://www.vidaventures.com/news",                          "depth": 1},
    # RTW Investments
    {"vc": "RTW Investments",           "url": "https://www.rtwfunds.com/insights",                          "depth": 1},
    # Polaris Partners
    {"vc": "Polaris Partners",          "url": "https://www.polarispartners.com/news",                       "depth": 1},
    # Bain Capital Life Sciences
    {"vc": "Bain Capital LS",           "url": "https://www.baincapitallifesciences.com/news",               "depth": 1},
    # Perceptive Advisors
    {"vc": "Perceptive Advisors",       "url": "https://www.perceptiveadvisors.com/news",                    "depth": 1},
    # Deerfield Management
    {"vc": "Deerfield Management",      "url": "https://www.deerfield.com/news",                             "depth": 1},
    # SR One (GSK spin-out)
    {"vc": "SR One",                    "url": "https://www.srone.com/news",                                 "depth": 1},
    # GV (Google Ventures) Life Sciences
    {"vc": "GV Life Sciences",          "url": "https://www.gv.com/portfolio",                               "depth": 0},
    # Andreessen Horowitz (a16z) Bio + Health
    {"vc": "a16z Bio",                  "url": "https://a16z.com/bio-health/",                               "depth": 1},
    {"vc": "a16z Bio",                  "url": "https://a16z.com/tag/bio/",                                  "depth": 1},
    # Bessemer Venture Partners Health
    {"vc": "Bessemer Health",           "url": "https://www.bvp.com/portfolio#healthcare",                   "depth": 0},
    {"vc": "Bessemer Health",           "url": "https://www.bvp.com/memos",                                  "depth": 1},
    # MPM BioImpact Capital
    {"vc": "MPM BioImpact",             "url": "https://www.mpmbio.com/news",                                "depth": 1},
    # Frazier Life Sciences
    {"vc": "Frazier Life Sciences",     "url": "https://www.frazierhealthcare.com/news",                     "depth": 1},
    # Longitude Capital
    {"vc": "Longitude Capital",         "url": "https://www.longitudecapital.com/portfolio",                 "depth": 0},
    # Pivotal Life Sciences
    {"vc": "Pivotal Life Sciences",     "url": "https://pivotallifesciences.com/portfolio",                  "depth": 0},
    # New Enterprise Associates (NEA) Healthcare
    {"vc": "NEA Healthcare",            "url": "https://www.nea.com/portfolio?sector=health",                "depth": 0},
    # Venrock Healthcare Capital
    {"vc": "Venrock Healthcare",        "url": "https://www.venrock.com/portfolio/?cat=healthcare",          "depth": 0},
    # Alexandria Venture Investments
    {"vc": "Alexandria Venture",        "url": "https://www.alexandrialaunchpads.com/portfolio",             "depth": 0},
    # Canaan Partners Health
    {"vc": "Canaan Partners Health",    "url": "https://www.canaan.com/portfolio/?focus=health",             "depth": 0},
    # Leaps by Bayer
    {"vc": "Leaps by Bayer",            "url": "https://leaps.bayer.com/news",                               "depth": 1},
    # Merck Global Health Innovation Fund
    {"vc": "Merck GHI Fund",            "url": "https://www.merckghifund.com/portfolio",                     "depth": 0},
    # AbbVie Ventures
    {"vc": "AbbVie Ventures",           "url": "https://abbvieventures.com/portfolio",                       "depth": 0},
    # Bristol Myers Squibb Ventures
    {"vc": "BMS Ventures",              "url": "https://www.bristolmyerssquibbventures.com/portfolio",       "depth": 0},
    # Amgen Ventures
    {"vc": "Amgen Ventures",            "url": "https://www.amgenventures.com/",                             "depth": 0},
    # Pfizer Ventures
    {"vc": "Pfizer Ventures",           "url": "https://www.pfizer.com/news/press-releases",                 "depth": 1},
    # J&J Innovation (Janssen)
    {"vc": "J&J Innovation",            "url": "https://www.jnj.com/latest-news/category/innovation",        "depth": 1},
    {"vc": "Janssen R&D",               "url": "https://www.janssen.com/press-releases",                     "depth": 1},
    # Boehringer Ingelheim Venture Fund
    {"vc": "BI Venture Fund",           "url": "https://www.boehringer-ingelheim.com/innovation/venture-fund", "depth": 0},
    # Lilly Ventures (Eli Lilly)
    {"vc": "Lilly Ventures",            "url": "https://www.lillyventures.com/portfolio",                    "depth": 0},
    # Novartis Venture Fund
    {"vc": "Novartis Venture Fund",     "url": "https://www.novartisventurefund.com/portfolio",              "depth": 0},
    # Roche Venture Fund
    {"vc": "Roche Venture Fund",        "url": "https://www.rocheventurefund.com/",                          "depth": 0},
    # MRL Ventures Fund (Merck R&D)
    {"vc": "MRL Ventures Fund",         "url": "https://www.mrlventuresfund.com/portfolio",                  "depth": 0},
    # HealthQuest Capital
    {"vc": "HealthQuest Capital",       "url": "https://hq.capital/insights",                                "depth": 1},
    # RiverVest Venture Partners
    {"vc": "RiverVest",                 "url": "https://www.rivervest.com/portfolio",                        "depth": 0},
    # Pappas Capital
    {"vc": "Pappas Capital",            "url": "https://www.pappascapital.com/portfolio",                    "depth": 0},
    # Sofinnova Investments (US, separate from EU)
    {"vc": "Sofinnova Investments",     "url": "https://sofinnova.com/portfolio/",                           "depth": 0},
    # Agent Capital
    {"vc": "Agent Capital",             "url": "https://www.agentcap.com/portfolio",                         "depth": 0},
    # Forbion Capital Partners (EU)
    {"vc": "Forbion",                   "url": "https://www.forbion.com/portfolio",                          "depth": 0},
    # Sofinnova Crossover I (late-stage)
    {"vc": "Sofinnova Crossover",       "url": "https://www.sofinnova.com/crossover/portfolio",              "depth": 0},
    # Index Ventures Life Sciences
    {"vc": "Index Ventures Life Sci",   "url": "https://www.indexventures.com/companies/?sector=life-sciences", "depth": 0},
    # Advent Life Sciences (EU)
    {"vc": "Advent Life Sciences",      "url": "https://www.adventlifesciences.com/portfolio",               "depth": 0},
    # EQT Life Sciences (formerly LSP)
    {"vc": "EQT Life Sciences",         "url": "https://eqtgroup.com/investments/life-sciences/",            "depth": 0},

    # ═══════════════════════════════════════════════════════════════════════
    # BIOTECH COMPANY IR PAGES  —  high-quality pipeline slide PDFs
    # (Sorted by innovation tier / data richness)
    # Note: Q4CDN / JavaScript-rendered IR pages omitted — no PDF links in raw HTML.
    # These companies are covered by EDGAR slide downloads (SlideDownloader) and
    # their press-release RSS feeds below.
    # ═══════════════════════════════════════════════════════════════════════

    # ── Company press-release RSS feeds (static XML, no JS needed) ───────────
    # Gene editing / RNA / cell therapy
    {"vc": "Beam Therapeutics",         "url": "https://ir.beamtx.com/rss/news-releases",                  "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Intellia Therapeutics",     "url": "https://ir.intelliatx.com/rss/news-releases",              "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Editas Medicine",           "url": "https://ir.editasmedicine.com/rss/news-releases",           "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Sangamo Therapeutics",      "url": "https://investor.sangamo.com/rss/news-releases",           "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Alnylam Pharmaceuticals",   "url": "https://investors.alnylam.com/rss/news-releases",          "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Ionis Pharmaceuticals",     "url": "https://ir.ionispharma.com/rss/news-releases",             "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Arrowhead Research",        "url": "https://ir.arrowheadpharma.com/rss/news-releases",         "depth": 0, "rss": True, "subdir": "startup"},
    # Oncology precision medicine
    {"vc": "Blueprint Medicines",       "url": "https://ir.blueprintmedicines.com/rss/news-releases",      "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Relay Therapeutics",        "url": "https://ir.relaytx.com/rss/news-releases",                 "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Arvinas",                   "url": "https://ir.arvinas.com/rss/news-releases",                 "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Kymera Therapeutics",       "url": "https://ir.kymeratherapeutics.com/rss/news-releases",      "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "C4 Therapeutics",           "url": "https://ir.c4therapeutics.com/rss/news-releases",          "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Nurix Therapeutics",        "url": "https://ir.nurixtx.com/rss/news-releases",                 "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Recursion Pharmaceuticals", "url": "https://ir.recursion.com/rss/news-releases",               "depth": 0, "rss": True, "subdir": "startup"},
    # Rare disease / genetic
    {"vc": "Ultragenyx",                "url": "https://ir.ultragenyx.com/rss/news-releases",              "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Sarepta Therapeutics",      "url": "https://investorrelations.sarepta.com/rss/news-releases",  "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Agios Pharmaceuticals",     "url": "https://investor.agios.com/rss/news-releases",             "depth": 0, "rss": True, "subdir": "startup"},
    # Vaccines / mRNA
    {"vc": "Moderna",                   "url": "https://investors.modernatx.com/rss/news-releases",        "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "BioNTech",                  "url": "https://investors.biontech.de/rss/news-releases",          "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Novavax",                   "url": "https://ir.novavax.com/rss/news-releases",                 "depth": 0, "rss": True, "subdir": "startup"},
    # Large-cap biotech
    {"vc": "Regeneron",                 "url": "https://investor.regeneron.com/rss/news-releases",         "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Vertex Pharmaceuticals",    "url": "https://investors.vrtx.com/rss/news-releases",             "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Biogen",                    "url": "https://investors.biogen.com/rss/news-releases",           "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Gilead Sciences",           "url": "https://investors.gilead.com/rss/news-releases",           "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Incyte",                    "url": "https://investor.incyte.com/rss/news-releases",            "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Biohaven",                  "url": "https://ir.biohavenpharma.com/rss/news-releases",          "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Karuna Therapeutics",       "url": "https://ir.karunatx.com/rss/news-releases",               "depth": 0, "rss": True, "subdir": "startup"},
    # AI-driven
    {"vc": "Exscientia",                "url": "https://ir.exscientia.ai/rss/news-releases",               "depth": 0, "rss": True, "subdir": "startup"},
    # ── Startup / biotech news aggregators via PR Newswire / GlobeNewswire ──
    {"vc": "PR Newswire Biotech",       "url": "https://www.prnewswire.com/rss/news-releases-list.rss?category=BIOTECHNOLOGY", "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "GlobeNewswire Pharma",      "url": "https://www.globenewswire.com/RssFeed/industry/9100-Pharmaceuticals", "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "GlobeNewswire Biotech",     "url": "https://www.globenewswire.com/RssFeed/industry/9101-Biotechnology", "depth": 0, "rss": True, "subdir": "startup"},
    {"vc": "Business Wire Biotech",     "url": "https://feed.businesswire.com/rss/home/?rss=G22&rssid=20642", "depth": 0, "rss": True, "subdir": "startup"},

    # ═══════════════════════════════════════════════════════════════════════
    # HEALTHCARE CONFERENCE PRESENTATION ARCHIVES
    # ═══════════════════════════════════════════════════════════════════════
    # JPMorgan Healthcare Conference (January — most important annual event)
    {"vc": "JPM Healthcare Conf",       "url": "https://www.jpmorgan.com/about-us/ib/jpm-healthcare-conference", "depth": 1, "subdir": "conference"},
    {"vc": "JPM HC Coverage FierceBio", "url": "https://www.fiercebiotech.com/jp-morgan-healthcare-conference",  "depth": 1, "subdir": "conference"},
    # BIO International Convention
    {"vc": "BIO Convention",            "url": "https://www.bio.org/events/bio-international-convention",       "depth": 1, "subdir": "conference"},
    # ASCO Annual Meeting
    {"vc": "ASCO Annual Meeting",       "url": "https://www.asco.org/meetings-education/meetings/annual-meeting", "depth": 1, "subdir": "conference"},
    # ASH Annual Meeting
    {"vc": "ASH Annual Meeting",        "url": "https://www.hematology.org/meetings/annual-meeting",            "depth": 1, "subdir": "conference"},
    # AACR Annual Meeting
    {"vc": "AACR Annual Meeting",       "url": "https://www.aacr.org/meeting/aacr-annual-meeting-2025/",        "depth": 1, "subdir": "conference"},
    # ESMO Congress
    {"vc": "ESMO Congress",             "url": "https://www.esmo.org/meeting-calendar/esmo-congress",           "depth": 1, "subdir": "conference"},

    # ═══════════════════════════════════════════════════════════════════════
    # INDUSTRY NEWS RSS FEEDS  —  structured, no JS rendering needed
    # ═══════════════════════════════════════════════════════════════════════
    {"vc": "FierceBiotech",             "url": "https://www.fiercebiotech.com/rss/xml",                    "depth": 0, "rss": True},
    {"vc": "BioSpace",                  "url": "https://www.biospace.com/rss/news",                        "depth": 0, "rss": True},
    {"vc": "Endpoints News",            "url": "https://endpts.com/feed/",                                 "depth": 0, "rss": True},
    {"vc": "STAT News",                 "url": "https://www.statnews.com/feed/",                           "depth": 0, "rss": True},
    {"vc": "BioPharma Dive",            "url": "https://www.biopharmadive.com/feeds/news/",                "depth": 0, "rss": True},
    {"vc": "Drug Discovery Today",      "url": "https://www.drugdiscoverytoday.com/rss/news.xml",          "depth": 0, "rss": True},
    {"vc": "Clinical Trials Arena",     "url": "https://www.clinicaltrialsarena.com/feed/",                "depth": 0, "rss": True},

    # ═══════════════════════════════════════════════════════════════════════
    # VC / THOUGHT-LEADER BLOGS
    # ═══════════════════════════════════════════════════════════════════════
    {"vc": "a16z Bio Blog",             "url": "https://a16z.com/tag/bio/",                                 "depth": 1},
    {"vc": "Atlas LifeSciVC Blog",      "url": "https://lifescivc.com/",                                   "depth": 1},
    {"vc": "ARCH Blog (Medium)",        "url": "https://medium.com/arch-venture-partners",                  "depth": 1},
    {"vc": "Flagship Blog",             "url": "https://www.flagshippioneering.com/stories",               "depth": 1},
    {"vc": "RA Capital Blog",           "url": "https://www.racap.com/news",                               "depth": 1},
    {"vc": "BIO CEO Perspectives",      "url": "https://www.bio.org/blogs",                                "depth": 1},
]

# ── Patterns for finding slide / PDF / PPTX links ────────────────────────────

_PDF_RE   = re.compile(r"\.pdf(\?[^\"']*)?$",  re.IGNORECASE)
_PPTX_RE  = re.compile(r"\.pptx?(\?[^\"']*)?$", re.IGNORECASE)
_SLIDE_HOSTS = re.compile(
    r"(slideshare\.net|speakerdeck\.com|ir\..*\.com|investors\."
    r"|sec\.gov/Archives|s3\.amazonaws\.com|cdn\.|cloudfront\.net"
    r"|q4cdn\.com|q4web\.com|storage\.googleapis\.com"
    r"|content\.irwebpage\.com|d1io3yog0oux5\.cloudfront"
    r"|events\.q4inc\.com|event\.on24\.com)",
    re.IGNORECASE,
)
_SKIP_URL_RE = re.compile(
    r"(twitter\.com|linkedin\.com|facebook\.com|instagram\.com|youtube\.com"
    r"|mailto:|javascript:|#|\.css|\.js|\.png|\.jpg|\.gif|\.svg|\.ico"
    r"|\.xml|\.zip|\.xlsx|\.docx|\.mp4|\.mp3|\.wav)",
    re.IGNORECASE,
)
_PRESENTATION_RE = re.compile(
    r"(presentation|slide|pipeline|investor.day|annual.report|corporate.overview"
    r"|fact.sheet|webcast|conference|data.package|R&D.day|science.day"
    r"|corporate.deck|pitch.deck|business.overview|clinical.update"
    r"|pipeline.update|company.overview|investor.presentation)",
    re.IGNORECASE,
)

# ── NLP patterns ──────────────────────────────────────────────────────────────

_STAGE_RE = re.compile(
    r"\bphase\s*([123i]+|one|two|three)\b|preclinical\b|pivotal\b|first.in.human\b",
    re.IGNORECASE,
)
_STAGE_MAP = {
    "1": "phase1", "i": "phase1", "one": "phase1",
    "2": "phase2", "ii": "phase2", "two": "phase2",
    "3": "phase3", "iii": "phase3", "three": "phase3",
}
_DECISION_RE = re.compile(
    r"\b(discontinu|terminat|halted|failed|no.go|advance|initiat|"
    r"enroll|positive|approved|NDA|BLA|milestone)\b",
    re.IGNORECASE,
)
_INDICATION_RE = re.compile(
    r"\b(cancer|carcinoma|oncology|leukemia|lymphoma|melanoma|glioblastoma|"
    r"rare disease|autoimmune|rheumatoid|lupus|crohn|colitis|"
    r"neurology|alzheimer|parkinson|ALS|multiple sclerosis|"
    r"cardiovascular|heart failure|atherosclerosis|hypertension|"
    r"metabolic|diabetes|obesity|NASH|NAFLD|"
    r"infectious|HIV|hepatitis|influenza|COVID|SARS|"
    r"inflammation|psoriasis|atopic dermatitis|"
    r"hematology|sickle cell|hemophilia|gene therapy|CNS|psychiatric)\b",
    re.IGNORECASE,
)
_MECHANISM_RE = re.compile(
    r"\b(antibody|monoclonal|bispecific|ADC|antibody.drug conjugate|"
    r"small molecule|inhibitor|kinase|checkpoint|PD.?1|PD.?L1|CTLA.?4|"
    r"cell therapy|CAR.?T|T.?cell|NK cell|"
    r"gene therapy|CRISPR|AAV|lentiviral|"
    r"RNA|siRNA|mRNA|antisense|oligonucleotide|"
    r"enzyme|protein|peptide|vaccine)\b",
    re.IGNORECASE,
)
_INVEST_RE = re.compile(r"\$\s*(\d[\d,.]*)\s*(million|M|billion|B)\b", re.IGNORECASE)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _extract_stage(text: str) -> str:
    m = _STAGE_RE.search(text)
    if not m:
        return "unknown"
    raw = m.group(0).lower()
    if any(x in raw for x in ("preclinical", "first-in-human", "first in human")):
        return "preclinical"
    if "pivotal" in raw:
        return "phase3"
    grp = (m.group(1) or "").lower()
    return _STAGE_MAP.get(grp, "unknown")


def _extract_decision(text: str) -> str:
    lower = text.lower()
    if any(w in lower for w in ("discontinu", "terminat", "fail", "halt", "no-go", "pulled")):
        return "no-go"
    if any(w in lower for w in ("advance", "initiat", "enrol", "positive", "approved", "nda", "bla")):
        return "go"
    return "undecided"


def _extract_indication(text: str) -> str:
    m = _INDICATION_RE.search(text)
    return m.group(0).lower() if m else "unknown"


def _extract_mechanism(text: str) -> str:
    m = _MECHANISM_RE.search(text)
    return m.group(0).lower() if m else "unknown"


def _extract_investment(text: str) -> float:
    best = 0.0
    for m in _INVEST_RE.finditer(text):
        amt = float(m.group(1).replace(",", ""))
        mult = 1_000_000 if m.group(2).lower() in ("million", "m") else 1_000_000_000
        val = amt * mult
        if val > best:
            best = val
    return best


def _htm_to_text(html: bytes | str) -> str:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "nav", "header", "footer", "noscript"]):
        tag.decompose()
    return soup.get_text(separator=" ", strip=True)


def _pdf_to_text(raw: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            pages = [p.extract_text() or "" for p in pdf.pages[:60]]
        return "\n".join(pages)
    except Exception as exc:
        logger.debug("pdfplumber failed: %s", exc)
        return ""


def _pptx_to_text(raw: bytes) -> str:
    """Extract all text from a PowerPoint (.pptx) file."""
    try:
        from pptx import Presentation  # python-pptx
        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as exc:
        logger.debug("pptx extraction failed: %s", exc)
        return ""


def _url_to_filename(url: str) -> str:
    h = hashlib.md5(url.encode()).hexdigest()[:10]
    name = urlparse(url).path.rstrip("/").split("/")[-1] or "page"
    # sanitise
    name = re.sub(r"[^\w\-.]", "_", name)[:60]
    return f"{h}_{name}"


_CLINICAL_KW_RE = re.compile(
    r"\b(clinical|trial|biotech|pharma|drug|therapy|therapeutic|pipeline|"
    r"indication|efficacy|safety|patient|treatment|mechanism|compound|"
    r"IND|FDA|EMA|phase|preclinical|discontinued|approved|milestone)\b",
    re.IGNORECASE,
)


def _is_clinical_content(text: str, strict: bool = False) -> bool:
    """Filter pages with no clinical content.
    strict=True: requires both stage AND decision signal (for raw documents).
    strict=False: requires stage OR decision OR 3+ clinical keywords (for pages).
    """
    snippet = text[:6000]
    has_stage    = bool(_STAGE_RE.search(snippet))
    has_decision = bool(_DECISION_RE.search(snippet))
    if strict:
        return has_stage and has_decision
    if has_stage or has_decision:
        return True
    # fallback: at least 3 distinct clinical keyword hits
    return len(set(m.group(0).lower() for m in _CLINICAL_KW_RE.finditer(snippet))) >= 3


def _find_pdf_and_slide_links(html: str, base_url: str) -> list[str]:
    """Return all PDF / PPTX / presentation links found on the page."""
    soup = BeautifulSoup(html, "lxml")
    links: list[str] = []
    seen: set[str] = set()
    # Also check <iframe src=>, <embed src=>, <object data=> for embedded slides
    candidates = (
        [(tag, tag.get("href", "")) for tag in soup.find_all("a", href=True)]
        + [(tag, tag.get("src", ""))  for tag in soup.find_all(["iframe", "embed"], src=True)]
        + [(tag, tag.get("data", "")) for tag in soup.find_all("object",  data=True)]
    )
    for tag, href in candidates:
        href = (href or "").strip()
        if not href or _SKIP_URL_RE.search(href):
            continue
        full = urljoin(base_url, href)
        if full in seen:
            continue
        seen.add(full)
        is_pdf   = bool(_PDF_RE.search(full))
        is_pptx  = bool(_PPTX_RE.search(full))
        is_slide_host = bool(_SLIDE_HOSTS.search(full))
        link_text = tag.get_text(strip=True) if hasattr(tag, "get_text") else ""
        is_presentation = bool(_PRESENTATION_RE.search(link_text + " " + full))
        if is_pdf or is_pptx or is_slide_host or is_presentation:
            links.append(full)
    return links


# ── Main collector ────────────────────────────────────────────────────────────

class VCWebsiteCollector(BaseCollector):
    """
    Scrapes news/publications pages of top-tier bio VC firms, biotech IR pages,
    and healthcare conference archives.  Downloads PDFs and PPTX slide decks.

    Output saved to:  data/slides/<subdir>/<vc_slug>/
      subdir = "vc"          for VC firm pages (default)
      subdir = "startup"     for biotech company IR presentations
      subdir = "conference"  for healthcare conference archives

    Extracts clinical pipeline signals (stage, decision, indication, mechanism)
    via NLP on page / document text.
    """

    name = "vc_website"
    rate_limit_seconds = 1.0  # polite crawling

    def __init__(
        self,
        base_slides_dir: Path = Path("data/slides"),
        timeout: int = 10,
    ) -> None:
        super().__init__(timeout=timeout)
        self.base_slides_dir = Path(base_slides_dir)
        # Pre-create all sub-directories
        for sub in ("vc", "startup", "conference"):
            (self.base_slides_dir / sub).mkdir(parents=True, exist_ok=True)
        # Legacy attribute kept for backward compat
        self.slides_dir = self.base_slides_dir / "vc"

    def collect(self, max_records: int = 500) -> list[RawRecord]:
        records: list[RawRecord] = []
        per_vc = max(max_records // max(len(_VC_PAGES), 1), 2)

        for page_cfg in _VC_PAGES:
            if len(records) >= max_records:
                break
            vc     = page_cfg["vc"]
            url    = page_cfg["url"]
            depth  = page_cfg.get("depth", 1)
            is_rss = page_cfg.get("rss", False)
            subdir = page_cfg.get("subdir", "vc")
            out_dir = self.base_slides_dir / subdir

            logger.info("[vc_website/%s] %s → %s", subdir, vc, url)
            if is_rss:
                batch = self._collect_from_rss(vc, url, per_vc)
            else:
                batch = self._collect_from_page(vc, url, depth, per_vc, out_dir)
            records.extend(batch)
            logger.info("  → %d records (total so far: %d)", len(batch), len(records))

        logger.info("[vc_website] Total: %d records", len(records))
        return records[:max_records]

    # ── RSS feed collection ───────────────────────────────────────────────────

    def _collect_from_rss(self, vc: str, url: str, limit: int) -> list[RawRecord]:
        """Parse an RSS/Atom feed and extract clinical content from each item."""
        try:
            resp = self._get(url, accept_json=False)
        except Exception as exc:
            logger.warning("  ✗ RSS fetch failed %s: %s", url, exc)
            return []

        soup = BeautifulSoup(resp.text, "xml")
        items = soup.find_all("item") or soup.find_all("entry")
        records: list[RawRecord] = []

        for item in items:
            if len(records) >= limit:
                break
            title_tag   = item.find("title")
            link_tag    = item.find("link")
            desc_tag    = item.find("description") or item.find("summary") or item.find("content")

            title_text = title_tag.get_text(strip=True) if title_tag else ""
            link_url   = (link_tag.get_text(strip=True) if link_tag else "") or url
            desc_text  = desc_tag.get_text(" ", strip=True) if desc_tag else ""
            full_text  = f"{title_text} {desc_text}"

            if not _is_clinical_content(full_text, strict=False):
                continue

            rec = self._make_record(vc, link_url, url, title_text, full_text)
            if rec:
                records.append(rec)

        return records

    # ── per-page logic ────────────────────────────────────────────────────────

    def _collect_from_page(
        self, vc: str, url: str, depth: int, limit: int,
        out_dir: Path | None = None,
    ) -> list[RawRecord]:
        if out_dir is None:
            out_dir = self.slides_dir

        records: list[RawRecord] = []

        try:
            resp = self._get(url, accept_json=False)
        except Exception as exc:
            logger.warning("  ✗ Failed to fetch %s: %s", url, exc)
            return records

        html = resp.text

        # Extract text from the landing page itself
        page_text = _htm_to_text(html)
        if _is_clinical_content(page_text, strict=False):
            rec = self._make_record(vc, url, url, page_text[:4000], page_text)
            if rec:
                records.append(rec)

        if depth == 0:
            return records

        # Find and follow PDF / PPTX / slide links
        linked_urls = _find_pdf_and_slide_links(html, url)
        logger.debug("  Found %d linked docs on %s", len(linked_urls), url)

        for linked_url in linked_urls:
            if len(records) >= limit:
                break
            text = self._fetch_linked(vc, linked_url, out_dir)
            if not text or not _is_clinical_content(text, strict=True):
                continue
            rec = self._make_record(vc, linked_url, url, page_text[:500], text)
            if rec:
                records.append(rec)

        return records

    def _fetch_linked(
        self, vc: str, url: str, out_dir: Path | None = None
    ) -> str:
        """Download a linked PDF / PPTX / HTML; return extracted text."""
        if out_dir is None:
            out_dir = self.slides_dir
        vc_slug   = re.sub(r"[^\w]", "_", vc.lower())[:20]
        vc_dir    = out_dir / vc_slug
        vc_dir.mkdir(parents=True, exist_ok=True)

        filename = _url_to_filename(url)
        is_pdf   = bool(_PDF_RE.search(url))
        is_pptx  = bool(_PPTX_RE.search(url))
        if is_pdf:
            local_ext = ".pdf"
        elif is_pptx:
            local_ext = ".pptx"
        else:
            local_ext = ".html"
        local_path = vc_dir / (filename + local_ext)

        # Cache hit
        if local_path.exists() and local_path.stat().st_size > 200:
            logger.debug("  Cache hit: %s", local_path.name)
            raw = local_path.read_bytes()
        else:
            try:
                resp = self._get(url, accept_json=False)
                raw = resp.content
                local_path.write_bytes(raw)
                logger.info("  ✓ Saved %s (%d KB) [%s]", local_path.name, len(raw)//1024, vc)
            except Exception as exc:
                logger.debug("  Download failed %s: %s", url, exc)
                return ""

        if is_pdf:
            return _pdf_to_text(raw)
        if is_pptx:
            return _pptx_to_text(raw)
        return _htm_to_text(raw)

    def _make_record(
        self,
        vc: str,
        url: str,
        referrer: str,
        page_snippet: str,
        full_text: str,
    ) -> RawRecord | None:
        text_slice = full_text[:8000]

        stage      = _extract_stage(text_slice)
        decision   = _extract_decision(text_slice)
        indication = _extract_indication(text_slice)
        mechanism  = _extract_mechanism(text_slice)
        investment = _extract_investment(text_slice)

        if stage == "unknown" and decision == "undecided":
            return None

        source_id = hashlib.md5(url.encode()).hexdigest()[:16]
        parsed    = urlparse(url)
        domain    = parsed.netloc.replace("www.", "")

        # Use page_snippet as title if it looks like an article headline (<120 chars)
        if page_snippet and len(page_snippet) < 120 and page_snippet.count("\n") == 0:
            title = f"{vc} — {page_snippet}"
        else:
            title = f"{vc} — {domain} ({stage})"

        outcome = "ongoing"
        if decision == "no-go":
            s = stage if stage not in ("unknown", "") else "p2"
            outcome = f"discontinued_{s}"
        elif decision == "go" and stage == "phase3":
            if "approved" in full_text.lower():
                outcome = "approved"

        return RawRecord(
            source         = self.name,
            source_id      = source_id,
            url            = url,
            title          = title[:200],
            indication     = indication,
            mechanism      = mechanism,
            clinical_stage = stage,
            decision       = decision,
            outcome        = outcome,
            investment_usd = investment,
            raw_text       = full_text[:4000],
            extra          = {
                "vc":        vc,
                "referrer":  referrer,
                "domain":    domain,
            },
        )
